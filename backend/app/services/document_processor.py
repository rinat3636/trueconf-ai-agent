import os
import re
import hashlib
from typing import List, Tuple, Optional

import pandas as pd
from docx import Document as DocxDocument
from PyPDF2 import PdfReader

from app.core.config import settings


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def extract_text_from_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    texts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            texts.append(text)
    result = "\n\n".join(texts)
    if result.strip():
        return result
    # Fallback: OCR for scanned PDFs
    try:
        from pdf2image import convert_from_path
        import pytesseract
        images = convert_from_path(file_path, dpi=200)
        ocr_texts = []
        for img in images:
            ocr_text = pytesseract.image_to_string(img, lang='rus+eng')
            if ocr_text.strip():
                ocr_texts.append(ocr_text)
        return "\n\n".join(ocr_texts)
    except Exception:
        return result


def extract_text_from_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    texts = []
    for para in doc.paragraphs:
        if para.text.strip():
            texts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_data:
                texts.append(" | ".join(row_data))
    return "\n\n".join(texts)


def extract_text_from_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"--- Слайд {slide_num} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_data:
                        slide_texts.append(" | ".join(row_data))
        if len(slide_texts) > 1:
            texts.append("\n".join(slide_texts))
    return "\n\n".join(texts)


def extract_text_from_xlsx(file_path: str) -> str:
    texts = []
    try:
        xls = pd.ExcelFile(file_path, engine="openpyxl")
    except Exception:
        xls = pd.ExcelFile(file_path, engine="xlrd")

    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name, header=None)
        texts.append(f"=== {sheet_name} ===")
        for _, row in df.iterrows():
            row_values = [str(v) for v in row.values if pd.notna(v) and str(v).strip()]
            if row_values:
                texts.append(" | ".join(row_values))
    return "\n".join(texts)


def extract_text_from_csv(file_path: str) -> str:
    texts = []
    for encoding in ["utf-8", "cp1251", "latin-1"]:
        try:
            df = pd.read_csv(file_path, encoding=encoding, header=None)
            for _, row in df.iterrows():
                row_values = [str(v) for v in row.values if pd.notna(v) and str(v).strip()]
                if row_values:
                    texts.append(" | ".join(row_values))
            return "\n".join(texts)
        except UnicodeDecodeError:
            continue
    return ""


def extract_text_from_txt(file_path: str) -> str:
    for encoding in ["utf-8", "cp1251", "latin-1"]:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    return ""


def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    extractors = {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".pptx": extract_text_from_pptx,
        ".xlsx": extract_text_from_xlsx,
        ".xls": extract_text_from_xlsx,
        ".csv": extract_text_from_csv,
        ".txt": extract_text_from_txt,
    }
    extractor = extractors.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported file type: {ext}")
    return extractor(file_path)


def compute_file_checksum(file_path: str) -> str:
    sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            sha256.update(chunk)
    return sha256.hexdigest()


# ---------------------------------------------------------------------------
# Smart chunking strategies (from ARCHITECTURE.md 4.3)
# ---------------------------------------------------------------------------

def chunk_text_semantic(
    text: str,
    chunk_size: int = 800,
    overlap: int = 100,
) -> List[str]:
    """Semantic chunking for DOCX/PDF text documents.
    Split by paragraph boundaries, respecting overlap."""
    if not text:
        return []

    separators = ["\n\n", "\n", ". ", "; "]
    chunks = []
    current_chunk = ""

    paragraphs = text.split("\n\n")
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        if len(current_chunk) + len(para) + 2 <= chunk_size:
            current_chunk += ("\n\n" + para) if current_chunk else para
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            if len(para) > chunk_size:
                sub_chunks = _split_long_paragraph(para, chunk_size, overlap)
                chunks.extend(sub_chunks[:-1])
                current_chunk = sub_chunks[-1] if sub_chunks else ""
            else:
                overlap_text = current_chunk[-overlap:] if current_chunk and overlap else ""
                current_chunk = (overlap_text + "\n\n" + para).strip() if overlap_text else para

    if current_chunk.strip():
        chunks.append(current_chunk.strip())

    return [c for c in chunks if len(c) > 20]


def _split_long_paragraph(text: str, chunk_size: int, overlap: int) -> List[str]:
    """Split a long paragraph into overlapping chunks by sentence boundaries."""
    sentences = re.split(r'(?<=[.!?;])\s+', text)
    chunks = []
    current = ""

    for sentence in sentences:
        if len(current) + len(sentence) + 1 <= chunk_size:
            current += (" " + sentence) if current else sentence
        else:
            if current:
                chunks.append(current.strip())
            overlap_words = current.split()[-overlap // 5:] if current else []
            current = " ".join(overlap_words) + " " + sentence if overlap_words else sentence

    if current.strip():
        chunks.append(current.strip())

    return chunks


def chunk_text_section_based(
    text: str,
    max_chunk_size: int = 1200,
) -> List[str]:
    """Section-based chunking for policy/regulation documents.
    Splits by numbered sections (1., 2., 1.1, etc.) and preserves headers."""
    if not text:
        return []

    section_pattern = re.compile(r'^(\d+(?:\.\d+)*\.?\s+)', re.MULTILINE)
    parts = section_pattern.split(text)

    sections = []
    current_header = ""
    current_body = ""

    i = 0
    while i < len(parts):
        part = parts[i].strip()
        if section_pattern.match(part + " "):
            if current_body:
                sections.append((current_header + current_body).strip())
            current_header = part
            i += 1
            if i < len(parts):
                current_body = parts[i]
            i += 1
        else:
            current_body += "\n" + part if current_body else part
            i += 1

    if current_body:
        sections.append((current_header + current_body).strip())

    chunks = []
    for section in sections:
        if len(section) <= max_chunk_size:
            chunks.append(section)
        else:
            sub_chunks = chunk_text_semantic(section, max_chunk_size, overlap=100)
            chunks.extend(sub_chunks)

    return [c for c in chunks if len(c) > 20]


def chunk_product_rows(
    file_path: str,
) -> List[str]:
    """Row-based chunking for product catalog XLS/XLSX files.
    Each product row becomes its own chunk using structured template."""
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".xls":
            xls = pd.ExcelFile(file_path, engine="xlrd")
        else:
            xls = pd.ExcelFile(file_path, engine="openpyxl")
    except Exception:
        return []

    chunks = []
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name)
        if df.empty:
            continue

        columns_lower = [str(c).lower().strip() for c in df.columns]

        name_col = _find_column(columns_lower, ["наименование", "название", "продукт", "товар", "name"])
        if name_col is None:
            continue

        for _, row in df.iterrows():
            name_val = str(row.iloc[name_col]).strip() if pd.notna(row.iloc[name_col]) else ""
            if not name_val or name_val == "nan":
                continue

            parts = [f"Продукт: {name_val}"]

            for label, keywords in [
                ("Артикул", ["артикул", "article", "код"]),
                ("Состав", ["состав", "composition"]),
                ("БЗМЖ/ЗМЖ", ["бзмж", "змж", "zmj"]),
                ("Срок годности", ["срок", "годност", "shelf"]),
                ("Штрихкод", ["штрих", "barcode", "ean"]),
                ("Вес", ["вес", "масса", "weight", "нетто"]),
                ("Упаковка", ["упаковк", "коробк", "блок", "pack"]),
                ("Паллетизация", ["паллет", "pallet"]),
                ("Пищевая ценность", ["пищев", "бжу", "калор", "энерг"]),
                ("Декларация", ["деклар", "гост", "ту"]),
                ("Категория", ["категори", "group", "тип"]),
                ("НДС", ["ндс", "nds", "vat"]),
            ]:
                col_idx = _find_column(columns_lower, keywords)
                if col_idx is not None:
                    val = str(row.iloc[col_idx]).strip() if pd.notna(row.iloc[col_idx]) else ""
                    if val and val != "nan":
                        parts.append(f"{label}: {val}")

            if len(parts) > 1:
                chunks.append(". ".join(parts) + ".")

    return chunks


def chunk_contacts(text: str) -> List[str]:
    """Row-based chunking for contact lists."""
    if not text:
        return []

    chunks = []
    lines = text.split("\n")
    for line in lines:
        line = line.strip()
        if not line or len(line) < 5:
            continue
        has_contact_info = any(kw in line.lower() for kw in [
            "@", "тел", "tel", "phone", "email", "почт",
            "внутр", "ext", "доб",
        ])
        has_name = bool(re.search(r'[А-ЯЁ][а-яё]+\s+[А-ЯЁ][а-яё]+', line))
        if has_contact_info or has_name:
            chunks.append(line)

    return chunks


def _find_column(columns: List[str], keywords: List[str]) -> Optional[int]:
    """Find column index by keyword match."""
    for i, col in enumerate(columns):
        for kw in keywords:
            if kw in col:
                return i
    return None


def detect_document_category(filename: str, text: str) -> str:
    """Auto-detect document category based on filename and content."""
    filename_lower = filename.lower()
    text_lower = text[:3000].lower() if text else ""

    if any(kw in filename_lower for kw in ["каталог", "catalog", "прайс", "price"]):
        return "product_catalog"
    if any(kw in filename_lower for kw in ["логист", "logist", "палле", "pallet"]):
        return "logistics"
    if any(kw in filename_lower for kw in ["полити", "policy", "пдз", "регламент"]):
        return "policy"
    if any(kw in filename_lower for kw in ["коммерч", "commercial", "кп"]):
        return "commercial"
    if any(kw in filename_lower for kw in ["деклар", "сертиф", "гост"]):
        return "certification"
    if any(kw in filename_lower for kw in ["контакт", "телефон", "почт", "email"]):
        return "contacts"
    if any(kw in filename_lower for kw in ["продаж", "прибыл", "выручк", "отчет", "отчёт"]):
        return "sales_data"
    if any(kw in filename_lower for kw in ["обуч", "класс", "трейн", "train"]):
        return "product_knowledge"
    if any(kw in filename_lower for kw in ["метод", "книга продаж"]):
        return "sales_methodology"

    if "кредитн" in text_lower or "дебиторск" in text_lower or "стоп-лист" in text_lower:
        return "policy"
    if "артикул" in text_lower and "штрихкод" in text_lower:
        return "product_catalog"
    if any(kw in text_lower for kw in ["валовая прибыль", "рентабельность", "тоннаж"]):
        return "sales_data"
    if "@" in text_lower and ("тел" in text_lower or "ext" in text_lower):
        return "contacts"

    return "general"


def get_chunking_strategy(category: str, file_ext: str) -> str:
    """Determine chunking strategy based on category and file type."""
    if category in ("product_catalog", "logistics") and file_ext in (".xls", ".xlsx"):
        return "row_based_product"
    if category in ("contacts",):
        return "contacts"
    if category in ("policy", "debt_management"):
        return "section_based"
    if file_ext in (".xls", ".xlsx", ".csv"):
        return "row_based_generic"
    return "semantic"


def smart_chunk(
    file_path: str,
    text: str,
    category: str,
) -> List[str]:
    """Apply the correct chunking strategy based on document category and type."""
    ext = os.path.splitext(file_path)[1].lower()
    strategy = get_chunking_strategy(category, ext)

    if strategy == "row_based_product":
        chunks = chunk_product_rows(file_path)
        if chunks:
            return chunks
        return chunk_text_semantic(text, settings.CHUNK_SIZE, settings.CHUNK_OVERLAP)

    if strategy == "contacts":
        chunks = chunk_contacts(text)
        if chunks:
            return chunks
        return chunk_text_semantic(text, settings.CHUNK_SIZE, settings.CHUNK_OVERLAP)

    if strategy == "section_based":
        return chunk_text_section_based(text, max_chunk_size=1200)

    return chunk_text_semantic(text, settings.CHUNK_SIZE, settings.CHUNK_OVERLAP)


# Legacy alias
def chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
    return chunk_text_semantic(text, chunk_size, chunk_overlap)


# ---------------------------------------------------------------------------
# Sales report parsing
# ---------------------------------------------------------------------------

# Patterns for hierarchy detection in 1C sales reports
_FIO_PATTERN = re.compile(r'^[А-ЯЁ][а-яё]+ [А-ЯЁ][а-яё]+ [А-ЯЁ][а-яё]+$')
_PRODUCT_PATTERN = re.compile(
    r'(\d+\s*г[р.]?\s*[*×xх]\s*\d+'  # weight*pack: 60г*40
    r'|\d+\s*г\b'                      # weight: 60г
    r'|\(\s*[А-Яа-яA-Za-z]-?\d+\s*\)' # article: (А-7/1)
    r'|\d+\s*[мМ][лЛ]'                # ml: 100мл
    r'|\bшт\.?\b'                      # шт
    r'|Сибхолод|Айс-Групп|Айс-Гр\b|Русс\.?\s*холод|МТ-холод'  # manufacturers
    r'|\bАГ\b'                         # АГ = Айс-Групп abbreviation
    r'|Стак\.|Рожок|Эскимо|Трубочк|Торт|Вафл|Брикет|Пакет|Фруктов'  # product types
    r')',
    re.IGNORECASE,
)


def _is_manager_name(name: str) -> bool:
    return bool(_FIO_PATTERN.match(name))


def _is_product_name(name: str) -> bool:
    return bool(_PRODUCT_PATTERN.search(name))


def _read_excel(file_path: str):
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".xls":
            return pd.read_excel(file_path, engine="xlrd", header=None)
        else:
            return pd.read_excel(file_path, engine="openpyxl", header=None)
    except Exception:
        return pd.read_csv(file_path, header=None)


# Column name synonyms for flexible matching
_COLUMN_SYNONYMS = {
    "manager": ["менеджер", "тп", "торговый представитель", "manager", "rep", "сотрудник", "ответственный"],
    "client": ["клиент", "контрагент", "покупатель", "client", "customer", "точка", "организация"],
    "product": ["товар", "продукт", "номенклатура", "product", "sku", "артикул", "наименование товара"],
    "quantity": ["количество", "кол-во", "кол.", "qty", "quantity", "шт"],
    "tonnage": ["тоннаж", "вес", "масса", "кг", "tonnage", "weight"],
    "revenue": ["выручка", "стоимость", "сумма", "revenue", "оборот", "продажи", "сумма продаж"],
    "profit": ["прибыль", "валовая прибыль", "доход", "profit", "gross profit", "маржинальный доход"],
    "margin": ["маржа", "маржинальность", "рентабельность", "margin", "%", "наценка"],
}


def _detect_column_mapping(df, header_row: int) -> dict:
    """Auto-detect column mapping from header row using synonyms."""
    if header_row is None or header_row >= len(df):
        return {}

    header_vals = []
    row = df.iloc[header_row]
    for v in row.values:
        header_vals.append(str(v).strip().lower() if pd.notna(v) else "")

    mapping = {}
    for field, synonyms in _COLUMN_SYNONYMS.items():
        for col_idx, header in enumerate(header_vals):
            if not header:
                continue
            for syn in synonyms:
                if syn in header:
                    mapping[field] = col_idx
                    break
            if field in mapping:
                break

    return mapping


def _parse_1c_hierarchical(df, metadata: dict) -> Tuple[List[dict], dict]:
    """Original 1C hierarchical parser (ТП→Клиент→Продукт nested structure)."""
    records = []

    header_row = None
    for idx, row in df.iterrows():
        vals = [str(v).strip() for v in row.values if pd.notna(v)]
        combined = " ".join(vals).lower()
        if "количество" in combined and ("стоимость" in combined or "прибыль" in combined):
            header_row = idx
            break

    if header_row is None:
        return records, metadata

    current_manager = None
    current_client = None

    for idx in range(header_row + 3, len(df)):
        row = df.iloc[idx]
        values = [v for v in row.values]

        name_val = str(values[1]).strip() if len(values) > 1 and pd.notna(values[1]) else ""
        if not name_val:
            continue

        nums = []
        for v in values[2:]:
            if pd.notna(v):
                try:
                    nums.append(float(v))
                except (ValueError, TypeError):
                    nums.append(None)
            else:
                nums.append(None)

        has_numbers = any(n is not None for n in nums)
        if not has_numbers:
            continue

        quantity = nums[0] if len(nums) > 0 else None
        tonnage = nums[1] if len(nums) > 1 else None
        revenue = nums[2] if len(nums) > 2 else None
        profit = nums[3] if len(nums) > 3 else None
        margin = nums[4] if len(nums) > 4 else None

        is_total = "итого" in name_val.lower()
        if is_total:
            continue

        if _is_manager_name(name_val):
            current_manager = name_val
            current_client = None
            records.append({
                "manager_name": current_manager,
                "client_name": None,
                "product_name": None,
                "quantity": quantity,
                "tonnage": tonnage,
                "revenue": revenue,
                "profit": profit,
                "margin_pct": margin,
                "record_level": "manager",
            })
        elif current_manager and _is_product_name(name_val):
            if not current_client:
                current_client = "(без клиента)"
            records.append({
                "manager_name": current_manager,
                "client_name": current_client,
                "product_name": name_val,
                "quantity": quantity,
                "tonnage": tonnage,
                "revenue": revenue,
                "profit": profit,
                "margin_pct": margin,
                "record_level": "product",
            })
        elif current_manager:
            current_client = name_val
            records.append({
                "manager_name": current_manager,
                "client_name": current_client,
                "product_name": None,
                "quantity": quantity,
                "tonnage": tonnage,
                "revenue": revenue,
                "profit": profit,
                "margin_pct": margin,
                "record_level": "client",
            })

    return records, metadata


def _parse_flat_table(df) -> Tuple[List[dict], dict]:
    """Fallback flat-table parser: each row has manager + client + product columns."""
    metadata = {}

    # Try to find header row with known column names
    header_row = None
    for idx in range(min(20, len(df))):
        row_vals = [str(v).strip().lower() for v in df.iloc[idx].values if pd.notna(v)]
        combined = " ".join(row_vals)
        # Look for rows that have at least 2 recognizable column names
        found = 0
        for synonyms in _COLUMN_SYNONYMS.values():
            if any(s in combined for s in synonyms):
                found += 1
        if found >= 2:
            header_row = idx
            break

    if header_row is None:
        return [], metadata

    mapping = _detect_column_mapping(df, header_row)
    if not mapping:
        return [], metadata

    records = []
    managers_seen = {}
    clients_seen = {}

    for idx in range(header_row + 1, len(df)):
        row = df.iloc[idx]

        def _get(field):
            col = mapping.get(field)
            if col is not None and col < len(row) and pd.notna(row.iloc[col]):
                return row.iloc[col]
            return None

        manager_name = str(_get("manager")).strip() if _get("manager") else None
        client_name = str(_get("client")).strip() if _get("client") else None
        product_name = str(_get("product")).strip() if _get("product") else None

        if not manager_name or manager_name in ("nan", "None", ""):
            continue

        is_total = "итого" in (manager_name or "").lower() or "итого" in (client_name or "").lower()
        if is_total:
            continue

        def _num(field):
            v = _get(field)
            if v is None:
                return None
            try:
                return float(v)
            except (ValueError, TypeError):
                return None

        quantity = _num("quantity")
        tonnage = _num("tonnage")
        revenue = _num("revenue")
        profit = _num("profit")
        margin = _num("margin")

        # Aggregate manager-level if not seen
        if manager_name and manager_name not in managers_seen:
            managers_seen[manager_name] = {
                "manager_name": manager_name,
                "client_name": None,
                "product_name": None,
                "quantity": 0, "tonnage": 0, "revenue": 0, "profit": 0, "margin_pct": None,
                "record_level": "manager",
            }

        # Track manager totals
        mgr = managers_seen.get(manager_name)
        if mgr:
            mgr["quantity"] = (mgr["quantity"] or 0) + (quantity or 0)
            mgr["revenue"] = (mgr["revenue"] or 0) + (revenue or 0)
            mgr["profit"] = (mgr["profit"] or 0) + (profit or 0)

        # Client-level record
        client_key = f"{manager_name}|{client_name}" if client_name else None
        if client_name and client_key and client_key not in clients_seen:
            clients_seen[client_key] = {
                "manager_name": manager_name,
                "client_name": client_name,
                "product_name": None,
                "quantity": 0, "tonnage": 0, "revenue": 0, "profit": 0, "margin_pct": None,
                "record_level": "client",
            }

        if client_key and client_key in clients_seen:
            cl = clients_seen[client_key]
            cl["quantity"] = (cl["quantity"] or 0) + (quantity or 0)
            cl["revenue"] = (cl["revenue"] or 0) + (revenue or 0)
            cl["profit"] = (cl["profit"] or 0) + (profit or 0)

        # Product-level record (each row)
        if product_name and product_name not in ("nan", "None", ""):
            records.append({
                "manager_name": manager_name,
                "client_name": client_name or "(без клиента)",
                "product_name": product_name,
                "quantity": quantity,
                "tonnage": tonnage,
                "revenue": revenue,
                "profit": profit,
                "margin_pct": margin,
                "record_level": "product",
            })

    # Compute margins for aggregated records
    for mgr in managers_seen.values():
        if mgr["revenue"] and mgr["revenue"] > 0 and mgr["profit"] is not None:
            mgr["margin_pct"] = round(mgr["profit"] / mgr["revenue"] * 100, 2)

    for cl in clients_seen.values():
        if cl["revenue"] and cl["revenue"] > 0 and cl["profit"] is not None:
            cl["margin_pct"] = round(cl["profit"] / cl["revenue"] * 100, 2)

    all_records = list(managers_seen.values()) + list(clients_seen.values()) + records

    if all_records:
        manager_records = [r for r in all_records if r["record_level"] == "manager"]
        metadata["total_revenue"] = sum(r["revenue"] or 0 for r in manager_records)
        metadata["total_profit"] = sum(r["profit"] or 0 for r in manager_records)
        metadata["total_managers"] = len(manager_records)
        client_records = [r for r in all_records if r["record_level"] == "client"]
        metadata["total_clients"] = len(client_records)
        product_records = [r for r in all_records if r["record_level"] == "product"]
        unique_products = set(r["product_name"] for r in product_records if r["product_name"])
        metadata["total_skus"] = len(unique_products)
        metadata["parser"] = "flat_table"

    return all_records, metadata


def parse_sales_report(file_path: str) -> Tuple[List[dict], dict]:
    df = _read_excel(file_path)

    records = []
    metadata = {}

    # Extract period from header area
    for idx, row in df.iterrows():
        values = [v for v in row.values if pd.notna(v)]
        if not values:
            continue
        if "Период:" in str(row.values):
            for v in row.values:
                if pd.notna(v) and "Период:" in str(v):
                    metadata["period"] = str(v).replace("Период:", "").strip()

    # Try 1C hierarchical parser first
    records, metadata = _parse_1c_hierarchical(df, metadata)

    # Fallback to flat-table parser if 1C parser found nothing
    if not records:
        logger.info("1C parser found no records, trying flat-table parser for %s", file_path)
        records, flat_metadata = _parse_flat_table(df)
        metadata.update(flat_metadata)

    if records and "total_revenue" not in metadata:
        manager_records = [r for r in records if r["record_level"] == "manager"]
        metadata["total_revenue"] = sum(r["revenue"] or 0 for r in manager_records)
        metadata["total_profit"] = sum(r["profit"] or 0 for r in manager_records)
        metadata["total_managers"] = len(manager_records)
        client_records = [r for r in records if r["record_level"] == "client"]
        metadata["total_clients"] = len(client_records)
        product_records = [r for r in records if r["record_level"] == "product"]
        unique_products = set(r["product_name"] for r in product_records if r["product_name"])
        metadata["total_skus"] = len(unique_products)

    return records, metadata
